"""Index PDF and PowerPoint files for local semantic search."""

import argparse
import base64
import concurrent.futures
import hashlib
import io
import json
import os
import re
import shutil
import sqlite3
import struct
import subprocess
import sys
import tempfile
import time
import urllib.error
import urllib.request
from pathlib import Path

try:
    import sqlite_vec
except ImportError:
    sqlite_vec = None


DB_PATH = Path(os.environ.get("RAG_DB_PATH", ".rag/index.db"))
EMBED_PROVIDER = os.environ.get("RAG_EMBED_PROVIDER", "github").lower()
if EMBED_PROVIDER == "local":
    EMBED_MODEL = os.environ.get("RAG_EMBED_MODEL", "BAAI/bge-small-en-v1.5")
    EMBED_DIM = 384
    EMBED_BATCH_SIZE = 256
elif EMBED_PROVIDER == "github":
    EMBED_MODEL = os.environ.get("RAG_EMBED_MODEL", "text-embedding-3-small")
    EMBED_DIM = 1536
    EMBED_BATCH_SIZE = 100
else:
    raise RuntimeError(f"Unsupported embedding provider: {EMBED_PROVIDER}")
EMBED_URL = "https://models.inference.ai.azure.com/embeddings"
VISION_MODEL = os.environ.get("RAG_VISION_MODEL", "openai/gpt-4.1-nano")
VISION_URL = "https://models.github.ai/inference/chat/completions"
VISION_REQUEST_INTERVAL_SECONDS = 4.1
VISION_CONCURRENCY = int(os.environ.get("RAG_VISION_CONCURRENCY", "5"))
VISION_CACHE_DIR = Path(os.environ.get("RAG_VISION_CACHE_DIR", ".rag/vision-cache"))
MAX_REMOTE_RETRY_SECONDS = 60
SUPPORTED_EXTENSIONS = {".pdf", ".pptx"}
_local_embedder = None
MONTH_NUMBERS = {
    "jan": 1,
    "january": 1,
    "feb": 2,
    "february": 2,
    "mar": 3,
    "march": 3,
    "apr": 4,
    "april": 4,
    "may": 5,
    "jun": 6,
    "june": 6,
    "jul": 7,
    "july": 7,
    "aug": 8,
    "august": 8,
    "sep": 9,
    "sept": 9,
    "september": 9,
    "oct": 10,
    "october": 10,
    "nov": 11,
    "november": 11,
    "dec": 12,
    "december": 12,
}
MONTH_PATTERN = "|".join(sorted(MONTH_NUMBERS, key=len, reverse=True))


def get_db(db_path: Path = DB_PATH) -> sqlite3.Connection:
    db_path.parent.mkdir(parents=True, exist_ok=True)
    db = sqlite3.connect(db_path)
    db.row_factory = sqlite3.Row
    db.execute("PRAGMA foreign_keys = ON")
    db.execute("PRAGMA journal_mode = WAL")

    if sqlite_vec is None:
        raise RuntimeError("sqlite-vec is not installed. Run 'uv sync'.")
    db.enable_load_extension(True)
    sqlite_vec.load(db)
    db.enable_load_extension(False)

    db.executescript(
        """
        CREATE TABLE IF NOT EXISTS documents (
            id TEXT PRIMARY KEY,
            path TEXT NOT NULL UNIQUE,
            name TEXT NOT NULL,
            year INTEGER,
            month INTEGER,
            file_type TEXT NOT NULL,
            vision_model TEXT,
            content_hash TEXT NOT NULL,
            indexed_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP
        );
        CREATE TABLE IF NOT EXISTS chunks (
            id TEXT PRIMARY KEY,
            document_id TEXT NOT NULL REFERENCES documents(id) ON DELETE CASCADE,
            ordinal INTEGER NOT NULL,
            location TEXT NOT NULL,
            text TEXT NOT NULL,
            visual_description TEXT,
            visual_useful INTEGER NOT NULL DEFAULT 0
        );
        CREATE INDEX IF NOT EXISTS idx_chunks_document ON chunks(document_id);
        """
    )
    db.execute(
        f"""CREATE VIRTUAL TABLE IF NOT EXISTS vec_chunks USING vec0(
            chunk_id TEXT PRIMARY KEY,
            embedding float[{EMBED_DIM}]
        )"""
    )
    db.commit()
    return db


def sha256_bytes(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as source:
        for block in iter(lambda: source.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def stable_id(value: str) -> str:
    return hashlib.sha256(value.encode("utf-8")).hexdigest()


def normalize_text(text: str) -> str:
    return " ".join(text.replace("\x00", " ").split())


def chunk_text(text: str, max_chars: int = 2400, overlap_chars: int = 300) -> list[str]:
    """Split text into overlapping chunks, preferring sentence boundaries."""
    text = normalize_text(text)
    if not text:
        return []
    if max_chars <= overlap_chars:
        raise ValueError("max_chars must be greater than overlap_chars")

    chunks = []
    start = 0
    while start < len(text):
        end = min(start + max_chars, len(text))
        if end < len(text):
            boundary = text.rfind(". ", start + max_chars // 2, end)
            if boundary != -1:
                end = boundary + 1
        chunks.append(text[start:end].strip())
        if end == len(text):
            break
        start = end - overlap_chars
    return chunks


def extract_pdf(path: Path) -> list[tuple[str, str]]:
    from pypdf import PdfReader

    return [
        (f"page {page_number}", page.extract_text() or "")
        for page_number, page in enumerate(PdfReader(path).pages, start=1)
    ]


def extract_powerpoint(path: Path) -> list[tuple[str, str]]:
    from pptx import Presentation

    sections = []
    for slide_number, slide in enumerate(Presentation(path).slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    parts.append(" | ".join(cell.text for cell in row.cells))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                parts.append(f"Speaker notes: {notes}")
        sections.append((f"slide {slide_number}", "\n".join(parts)))
    return sections


def powerpoint_substantive_visuals(path: Path) -> list[bool]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    presentation = Presentation(path)
    slide_area = presentation.slide_width * presentation.slide_height
    visual_types = {
        MSO_SHAPE_TYPE.CHART,
        MSO_SHAPE_TYPE.MEDIA,
        MSO_SHAPE_TYPE.PICTURE,
        MSO_SHAPE_TYPE.LINKED_PICTURE,
    }
    results = []
    for slide in presentation.slides:
        substantive = False
        for shape in slide.shapes:
            if getattr(shape, "has_chart", False) or getattr(shape, "has_table", False):
                substantive = True
                break
            shape_area_ratio = (shape.width * shape.height) / slide_area
            if shape.shape_type in visual_types and shape_area_ratio >= 0.05:
                substantive = True
                break
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP and shape_area_ratio >= 0.10:
                substantive = True
                break
        results.append(substantive)
    return results


def windows_path(path: Path) -> str:
    resolved_path = path.resolve()
    if sys.platform == "win32":
        return str(resolved_path)
    if shutil.which("wslpath") is None:
        raise RuntimeError(
            "PowerPoint visual rendering requires Windows or WSL with Windows PowerPoint. "
            "Use --no-vision to skip visual rendering."
        )
    result = subprocess.run(
        ["wslpath", "-w", str(resolved_path)],
        capture_output=True,
        text=True,
        timeout=10,
        check=True,
    )
    return result.stdout.strip()


def render_powerpoint(path: Path, output_dir: Path) -> list[Path]:
    input_path = windows_path(path).replace("'", "''")
    output_path = windows_path(output_dir).replace("'", "''")
    command = f"""
$app = New-Object -ComObject PowerPoint.Application
try {{
    $deck = $app.Presentations.Open('{input_path}', $true, $true, $false)
    $deck.Export('{output_path}', 'JPG', 1280, 720)
    $deck.Close()
}} finally {{
    $app.Quit()
}}
"""
    encoded_command = base64.b64encode(command.encode("utf-16-le")).decode()
    subprocess.run(
        ["powershell.exe", "-NoProfile", "-NonInteractive", "-EncodedCommand", encoded_command],
        capture_output=True,
        text=True,
        timeout=180,
        check=True,
    )
    return rendered_slide_images(output_dir)


def rendered_slide_images(output_dir: Path) -> list[Path]:
    images_by_path = {}
    for pattern in ("*.JPG", "*.jpg"):
        for image in output_dir.glob(pattern):
            resolved_path = str(image.resolve())
            if sys.platform == "win32":
                resolved_path = resolved_path.casefold()
            images_by_path[resolved_path] = image
    images = list(images_by_path.values())
    return sorted(images, key=lambda image: int(re.search(r"(\d+)$", image.stem).group(1)))


def render_pdf(path: Path, output_dir: Path) -> list[Path]:
    import pypdfium2 as pdfium

    document = pdfium.PdfDocument(path)
    images = []
    try:
        for page_index in range(len(document)):
            image_path = output_dir / f"page-{page_index + 1}.jpg"
            bitmap = document[page_index].render(scale=1.5)
            bitmap.to_pil().convert("RGB").save(image_path, "JPEG", quality=80)
            images.append(image_path)
    finally:
        document.close()
    return images


def render_sections(path: Path, output_dir: Path) -> list[Path]:
    if path.suffix.lower() == ".pptx":
        return render_powerpoint(path, output_dir)
    if path.suffix.lower() == ".pdf":
        return render_pdf(path, output_dir)
    raise ValueError(f"Unsupported file type: {path.suffix}")


_last_vision_request = 0.0


def encode_image(path: Path) -> str:
    from PIL import Image

    with Image.open(path) as image:
        image.thumbnail((1280, 1280))
        buffer = io.BytesIO()
        image.convert("RGB").save(buffer, format="JPEG", quality=75, optimize=True)
    return "data:image/jpeg;base64," + base64.b64encode(buffer.getvalue()).decode()


def describe_image(path: Path, extracted_text: str) -> dict:
    global _last_vision_request
    content = [
        {
            "type": "text",
            "text": (
                "Describe this presentation slide for semantic retrieval. "
                "Focus on photographs, screenshots, charts, diagrams, UI state, status indicators, "
                "and spatial relationships that plain text extraction can miss. Include important "
                "visible labels or values. Set VISUAL_USEFUL to yes only when the image adds factual "
                "retrieval information beyond the extracted text supplied below. Logos, branding, "
                "backgrounds, formatting, ordinary text layout, and decorative elements do not count. "
                "Mandatory precedence rule: title/cover, agenda, section divider, intro, and text-only "
                "slides MUST be no even when they contain a photograph, logo, or decorative graphic. "
                "Return exactly two parts: a first line "
                "VISUAL_USEFUL: yes or VISUAL_USEFUL: no, followed by DESCRIPTION: and one concise "
                "factual paragraph.\n\nExtracted slide text:\n"
                + extracted_text[:6000]
            ),
        },
        {"type": "image_url", "image_url": {"url": encode_image(path)}},
    ]
    payload = {
        "model": VISION_MODEL,
        "messages": [{"role": "user", "content": content}],
        "max_tokens": 300,
        "temperature": 0,
    }
    data = json.dumps(payload).encode()
    headers = {
        "Authorization": f"Bearer {get_github_token()}",
        "Content-Type": "application/json",
    }
    for attempt in range(6):
        remaining_delay = VISION_REQUEST_INTERVAL_SECONDS - (
            time.monotonic() - _last_vision_request
        )
        if remaining_delay > 0:
            time.sleep(remaining_delay)
        request = urllib.request.Request(VISION_URL, data=data, headers=headers)
        try:
            with urllib.request.urlopen(request, timeout=90) as response:
                _last_vision_request = time.monotonic()
                body = json.loads(response.read())
                response_text = body["choices"][0]["message"]["content"].strip()
                useful_match = re.search(
                    r"(?im)^VISUAL_USEFUL:\s*(yes|no)\s*$", response_text
                )
                description_match = re.search(
                    r"(?is)DESCRIPTION:\s*(.+)$", response_text
                )
                if not useful_match or not description_match:
                    raise ValueError("Vision response did not include classification and description")
                return {
                    "useful": useful_match.group(1).lower() == "yes",
                    "description": description_match.group(1).strip(),
                }
        except urllib.error.HTTPError as error:
            if error.code != 429 or attempt == 5:
                raise
            retry_after = float(error.headers.get("Retry-After", max(10, 2**attempt)))
            if retry_after > MAX_REMOTE_RETRY_SECONDS:
                raise RuntimeError(
                    f"GitHub Models vision quota exhausted; retry after "
                    f"{int(retry_after)} seconds"
                ) from error
            time.sleep(retry_after)
        except (KeyError, TypeError, ValueError):
            if attempt == 5:
                raise
    raise RuntimeError("Vision request failed")


def describe_image_cached(path: Path, extracted_text: str) -> dict:
    cache_key = stable_id(
        f"{VISION_MODEL}:{sha256_bytes(path)}:{normalize_text(extracted_text)}"
    )
    cache_path = VISION_CACHE_DIR / f"{cache_key}.json"
    if cache_path.exists():
        return json.loads(cache_path.read_text(encoding="utf-8"))

    result = describe_image(path, extracted_text)
    cache_path.parent.mkdir(parents=True, exist_ok=True)
    temporary_path = cache_path.with_suffix(f".{os.getpid()}.tmp")
    temporary_path.write_text(json.dumps(result), encoding="utf-8")
    os.replace(temporary_path, cache_path)
    return result


def is_basic_slide(location: str, extracted_text: str) -> bool:
    section_match = re.fullmatch(
        r"(?:slide|page)\s+(\d+)", location, flags=re.IGNORECASE
    )
    if section_match and int(section_match.group(1)) == 1:
        return True

    normalized = normalize_text(extracted_text).lower()
    words = normalized.split()
    basic_phrases = (
        "agenda",
        "table of contents",
        "section overview",
        "thank you",
        "questions",
        "q&a",
        "appendix",
    )
    return len(words) <= 100 and any(
        normalized.startswith(phrase) for phrase in basic_phrases
    )


def add_visual_descriptions(
    path: Path,
    sections: list[tuple[str, str]],
    concurrency: int = VISION_CONCURRENCY,
    describe_mask: list[bool] | None = None,
) -> list[dict]:
    if concurrency < 1:
        raise ValueError("vision concurrency must be at least 1")
    if describe_mask is None:
        describe_mask = [not is_basic_slide(location, text) for location, text in sections]
    if len(describe_mask) != len(sections):
        raise ValueError("vision description mask must match section count")
    Path(".rag").mkdir(exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="render-", dir=".rag") as temporary_dir:
        images = render_sections(path, Path(temporary_dir).resolve())
        if len(images) != len(sections):
            raise RuntimeError(
                f"Rendered {len(images)} images for {len(sections)} sections in {path.name}"
            )
        candidates = [
            (index, image, sections[index][1])
            for index, image in enumerate(images)
            if describe_mask[index]
        ]
        with concurrent.futures.ThreadPoolExecutor(max_workers=concurrency) as executor:
            described_candidates = list(
                executor.map(
                    lambda item: describe_image_cached(item[1], item[2]),
                    candidates,
                )
            )
        visual_results = [
            {"description": None, "useful": False} for _ in sections
        ]
        for (index, _, _), visual_result in zip(candidates, described_candidates):
            visual_results[index] = visual_result
        return [
            {
                "location": location,
                "text": text,
                "visual_description": visual_result["description"],
                "visual_useful": visual_result["useful"]
                and not is_basic_slide(location, text),
            }
            for (location, text), visual_result in zip(sections, visual_results)
        ]


def extract_sections(
    path: Path,
    describe_images: bool = True,
    vision_concurrency: int = VISION_CONCURRENCY,
) -> list[dict]:
    if path.suffix.lower() == ".pdf":
        sections = extract_pdf(path)
        substantive_visuals = None
    elif path.suffix.lower() == ".pptx":
        sections = extract_powerpoint(path)
        substantive_visuals = powerpoint_substantive_visuals(path)
    else:
        raise ValueError(f"Unsupported file type: {path.suffix}")
    if describe_images:
        describe_mask = [
            not is_basic_slide(location, text)
            and (substantive_visuals[index] if substantive_visuals is not None else True)
            for index, (location, text) in enumerate(sections)
        ]
        described_sections = add_visual_descriptions(
            path, sections, vision_concurrency, describe_mask
        )
        return described_sections
    return [
        {
            "location": location,
            "text": text,
            "visual_description": None,
            "visual_useful": False,
        }
        for location, text in sections
    ]


def make_chunks(
    path: Path,
    root: Path,
    max_chars: int,
    overlap_chars: int,
    describe_images: bool,
    vision_concurrency: int,
) -> list[dict]:
    relative_path = path.relative_to(root).as_posix()
    chunks = []
    ordinal = 0
    for section in extract_sections(
        path, describe_images, vision_concurrency
    ):
        if path.suffix.lower() == ".pptx":
            text_chunks = [normalize_text(section["text"])]
        else:
            text_chunks = chunk_text(section["text"], max_chars, overlap_chars)
        if not text_chunks:
            text_chunks = [""]
        for section_index, section_chunk in enumerate(text_chunks):
            visual_description = (
                section["visual_description"] if section_index == 0 else None
            )
            visual_useful = bool(section["visual_useful"] and visual_description)
            embedding_parts = [section_chunk]
            if visual_useful:
                embedding_parts.append(visual_description)
            embedding_text = "\n\n".join(part for part in embedding_parts if part).strip()
            chunks.append(
                {
                    "id": stable_id(
                        f"{relative_path}:{ordinal}:{section_chunk}:{visual_description}:{visual_useful}"
                    ),
                    "ordinal": ordinal,
                    "location": section["location"],
                    "text": section_chunk,
                    "visual_description": visual_description,
                    "visual_useful": visual_useful,
                    "embedding_text": embedding_text,
                }
            )
            ordinal += 1
    return chunks


def get_github_token() -> str:
    environment_token = os.environ.get("GH_TOKEN") or os.environ.get("GITHUB_TOKEN")
    if environment_token:
        return environment_token

    result = subprocess.run(
        ["gh", "auth", "token"], capture_output=True, text=True, timeout=10, check=False
    )
    if result.returncode == 0 and result.stdout.strip():
        return result.stdout.strip()

    legacy_result = subprocess.run(
        ["gh", "auth", "status", "--show-token"],
        capture_output=True,
        text=True,
        timeout=10,
        check=False,
    )
    token_match = re.search(r"(?i)Token:\s*(\S+)", legacy_result.stdout + legacy_result.stderr)
    if token_match:
        return token_match.group(1)
    raise RuntimeError("No GitHub token available. Run 'gh auth login'.")


def get_embeddings(texts: list[str]) -> list[list[float]]:
    global _local_embedder
    if EMBED_PROVIDER == "local":
        if _local_embedder is None:
            from fastembed import TextEmbedding

            _local_embedder = TextEmbedding(model_name=EMBED_MODEL)
        return [
            embedding.tolist()
            for embedding in _local_embedder.embed(texts, batch_size=EMBED_BATCH_SIZE)
        ]

    payload = json.dumps({"input": texts, "model": EMBED_MODEL}).encode()
    headers = {
        "Authorization": f"Bearer {get_github_token()}",
        "Content-Type": "application/json",
    }
    for attempt in range(5):
        request = urllib.request.Request(EMBED_URL, data=payload, headers=headers)
        try:
            with urllib.request.urlopen(request, timeout=60) as response:
                body = json.loads(response.read())
                return [item["embedding"] for item in body["data"]]
        except urllib.error.HTTPError as error:
            if error.code != 429 or attempt == 4:
                raise
            retry_after = float(error.headers.get("Retry-After", 2**attempt))
            if retry_after > MAX_REMOTE_RETRY_SECONDS:
                raise RuntimeError(
                    f"GitHub Models embedding quota exhausted; retry after "
                    f"{int(retry_after)} seconds"
                ) from error
            time.sleep(retry_after)
    raise RuntimeError("Embedding request failed")


def serialize_f32(vector: list[float]) -> bytes:
    return struct.pack(f"{len(vector)}f", *vector)


def discover_documents(root: Path) -> list[Path]:
    return sorted(
        path
        for path in root.rglob("*")
        if path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS
    )


def parse_document_date(name: str) -> tuple[int | None, int | None]:
    month_then_year = re.search(
        rf"(?i)(?<![a-z])({MONTH_PATTERN})[\s._-]*(20\d{{2}})(?!\d)", name
    )
    if month_then_year:
        return int(month_then_year.group(2)), MONTH_NUMBERS[month_then_year.group(1).lower()]

    year_then_month = re.search(
        rf"(?i)(?<!\d)(20\d{{2}})[\s._-]*({MONTH_PATTERN})(?![a-z])", name
    )
    if year_then_month:
        return int(year_then_month.group(1)), MONTH_NUMBERS[year_then_month.group(2).lower()]

    numeric_date = re.search(r"(?<!\d)(20\d{2})[-_.](0?[1-9]|1[0-2])(?!\d)", name)
    if numeric_date:
        return int(numeric_date.group(1)), int(numeric_date.group(2))

    year = re.search(r"(?<!\d)(20\d{2})(?!\d)", name)
    return (int(year.group(1)), None) if year else (None, None)


def delete_document_chunks(db: sqlite3.Connection, document_id: str) -> None:
    old_chunk_ids = db.execute(
        "SELECT id FROM chunks WHERE document_id = ?", (document_id,)
    ).fetchall()
    for old_chunk in old_chunk_ids:
        db.execute("DELETE FROM vec_chunks WHERE chunk_id = ?", (old_chunk["id"],))
    db.execute("DELETE FROM chunks WHERE document_id = ?", (document_id,))


def upsert_document(
    db: sqlite3.Connection,
    document_id: str,
    relative_path: str,
    document_name: str,
    document_year: int | None,
    document_month: int | None,
    file_type: str,
    vision_model: str | None,
    file_hash: str,
) -> None:
    db.execute(
        """INSERT INTO documents
               (id, path, name, year, month, file_type, vision_model, content_hash)
             VALUES (?, ?, ?, ?, ?, ?, ?, ?)
          ON CONFLICT(id) DO UPDATE SET
            path = excluded.path,
            name = excluded.name,
            year = excluded.year,
            month = excluded.month,
            file_type = excluded.file_type,
            vision_model = excluded.vision_model,
            content_hash = excluded.content_hash,
            indexed_at = CURRENT_TIMESTAMP""",
        (
            document_id,
            relative_path,
            document_name,
            document_year,
            document_month,
            file_type,
            vision_model,
            file_hash,
        ),
    )


def index_document(
    db: sqlite3.Connection,
    path: Path,
    root: Path,
    max_chars: int,
    overlap_chars: int,
    describe_images: bool,
    vision_concurrency: int,
) -> dict:
    relative_path = path.relative_to(root).as_posix()
    document_name = path.stem
    document_year, document_month = parse_document_date(document_name)
    vision_model = VISION_MODEL if describe_images else None
    document_id = stable_id(relative_path)
    file_hash = sha256_bytes(path)
    existing = db.execute(
        "SELECT content_hash, vision_model FROM documents WHERE id = ?", (document_id,)
    ).fetchone()
    if (
        existing
        and existing["content_hash"] == file_hash
        and existing["vision_model"] == vision_model
    ):
        return {
            "name": document_name,
            "year": document_year,
            "month": document_month,
            "path": relative_path,
            "status": "unchanged",
            "chunks": 0,
        }

    chunks = make_chunks(
        path,
        root,
        max_chars,
        overlap_chars,
        describe_images,
        vision_concurrency,
    )
    if not chunks:
        delete_document_chunks(db, document_id)
        upsert_document(
            db,
            document_id,
            relative_path,
            document_name,
            document_year,
            document_month,
            path.suffix.lower().lstrip("."),
            vision_model,
            file_hash,
        )
        db.commit()
        return {
            "name": document_name,
            "year": document_year,
            "month": document_month,
            "path": relative_path,
            "status": "skipped",
            "chunks": 0,
            "reason": "no text",
        }

    embedded_chunks = []
    chunks_to_embed = [chunk for chunk in chunks if chunk["embedding_text"]]
    if not chunks_to_embed:
        delete_document_chunks(db, document_id)
        upsert_document(
            db,
            document_id,
            relative_path,
            document_name,
            document_year,
            document_month,
            path.suffix.lower().lstrip("."),
            vision_model,
            file_hash,
        )
        db.commit()
        return {
            "name": document_name,
            "year": document_year,
            "month": document_month,
            "path": relative_path,
            "status": "skipped",
            "chunks": 0,
            "reason": "no searchable text or useful visual description",
        }
    for offset in range(0, len(chunks_to_embed), EMBED_BATCH_SIZE):
        batch = chunks_to_embed[offset : offset + EMBED_BATCH_SIZE]
        vectors = get_embeddings([chunk["embedding_text"] for chunk in batch])
        embedded_chunks.extend(zip(batch, vectors))

    delete_document_chunks(db, document_id)
    upsert_document(
        db,
        document_id,
        relative_path,
        document_name,
        document_year,
        document_month,
        path.suffix.lower().lstrip("."),
        vision_model,
        file_hash,
    )

    for chunk in chunks:
        db.execute(
            """INSERT INTO chunks
                 (id, document_id, ordinal, location, text, visual_description, visual_useful)
               VALUES (?, ?, ?, ?, ?, ?, ?)""",
            (
                chunk["id"],
                document_id,
                chunk["ordinal"],
                chunk["location"],
                chunk["text"],
                chunk["visual_description"],
                int(chunk["visual_useful"]),
            ),
        )

    for chunk, vector in embedded_chunks:
        db.execute(
            "INSERT INTO vec_chunks (chunk_id, embedding) VALUES (?, ?)",
            (chunk["id"], serialize_f32(vector)),
        )
    db.commit()

    return {
        "name": document_name,
        "year": document_year,
        "month": document_month,
        "path": relative_path,
        "status": "indexed",
        "chunks": len(chunks),
    }


def prune_missing(db: sqlite3.Connection, discovered_paths: set[str]) -> int:
    removed = 0
    for document in db.execute("SELECT id, path FROM documents").fetchall():
        if document["path"] in discovered_paths:
            continue
        chunk_ids = db.execute(
            "SELECT id FROM chunks WHERE document_id = ?", (document["id"],)
        ).fetchall()
        for chunk in chunk_ids:
            db.execute("DELETE FROM vec_chunks WHERE chunk_id = ?", (chunk["id"],))
        db.execute("DELETE FROM documents WHERE id = ?", (document["id"],))
        removed += 1
    db.commit()
    return removed


def ingest(
    root: Path,
    db_path: Path = DB_PATH,
    max_chars: int = 2400,
    overlap_chars: int = 300,
    prune: bool = False,
    describe_images: bool = True,
    vision_concurrency: int = VISION_CONCURRENCY,
) -> dict:
    root = root.resolve()
    if not root.is_dir():
        raise ValueError(f"Document directory not found: {root}")
    db = get_db(db_path)
    documents = discover_documents(root)
    results = []
    for path in documents:
        try:
            results.append(
                index_document(
                    db,
                    path,
                    root,
                    max_chars,
                    overlap_chars,
                    describe_images,
                    vision_concurrency,
                )
            )
        except Exception as error:
            results.append(
                {
                    "path": path.relative_to(root).as_posix(),
                    "status": "error",
                    "error": str(error),
                }
            )
    discovered_paths = {path.relative_to(root).as_posix() for path in documents}
    removed = prune_missing(db, discovered_paths) if prune else 0
    return {
        "root": str(root),
        "found": len(documents),
        "indexed": sum(item["status"] == "indexed" for item in results),
        "unchanged": sum(item["status"] == "unchanged" for item in results),
        "skipped": sum(item["status"] == "skipped" for item in results),
        "errors": [item for item in results if item["status"] == "error"],
        "pruned": removed,
        "documents": results,
    }


def search(query: str, limit: int = 8, db_path: Path = DB_PATH) -> list[dict]:
    if not query.strip():
        raise ValueError("Query cannot be empty")
    db = get_db(db_path)
    query_vector = serialize_f32(get_embeddings([query])[0])
    rows = db.execute(
        """SELECT d.name, d.year, d.month, d.path, d.file_type,
              c.location, c.ordinal, c.text, c.visual_description,
              c.visual_useful, v.distance
           FROM vec_chunks AS v
           JOIN chunks AS c ON c.id = v.chunk_id
           JOIN documents AS d ON d.id = c.document_id
           WHERE v.embedding MATCH ? AND k = ?
           ORDER BY v.distance""",
        (query_vector, limit),
    ).fetchall()
    return [dict(row) for row in rows]


def stats(db_path: Path = DB_PATH) -> dict:
    db = get_db(db_path)
    return {
        "documents": db.execute("SELECT COUNT(*) FROM documents").fetchone()[0],
        "chunks": db.execute("SELECT COUNT(*) FROM chunks").fetchone()[0],
        "embeddings": db.execute("SELECT COUNT(*) FROM vec_chunks").fetchone()[0],
        "database": str(db_path),
        "embedding_provider": EMBED_PROVIDER,
        "model": EMBED_MODEL,
    }


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--db", type=Path, default=DB_PATH, help="SQLite index path")
    subparsers = parser.add_subparsers(dest="command", required=True)

    ingest_parser = subparsers.add_parser("ingest", help="Recursively index PDFs and PPTX files")
    ingest_parser.add_argument("directory", type=Path)
    ingest_parser.add_argument(
        "--max-chars", type=int, default=2400, help="Maximum PDF page chunk size"
    )
    ingest_parser.add_argument(
        "--overlap-chars", type=int, default=300, help="PDF page chunk overlap"
    )
    ingest_parser.add_argument("--prune", action="store_true", help="Remove records for missing files")
    ingest_parser.add_argument(
        "--no-vision",
        action="store_true",
        help="Skip slide/page rendering and visual descriptions",
    )
    ingest_parser.add_argument(
        "--vision-concurrency",
        type=int,
        default=VISION_CONCURRENCY,
        help=f"Concurrent vision requests (default: {VISION_CONCURRENCY})",
    )

    search_parser = subparsers.add_parser("search", help="Run semantic search")
    search_parser.add_argument("query")
    search_parser.add_argument("--limit", type=int, default=8)
    subparsers.add_parser("stats", help="Show index statistics")
    return parser


def main() -> None:
    args = build_parser().parse_args()
    if args.command == "ingest":
        result = ingest(
            args.directory,
            args.db,
            args.max_chars,
            args.overlap_chars,
            args.prune,
            not args.no_vision,
            args.vision_concurrency,
        )
    elif args.command == "search":
        result = search(args.query, args.limit, args.db)
    else:
        result = stats(args.db)
    print(json.dumps(result, indent=2))


if __name__ == "__main__":
    main()
